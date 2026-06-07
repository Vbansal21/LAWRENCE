"""Document converters for the UI bridge attachment pipeline.

Each converter receives a local Path (or URL string for web pages) and returns
extracted plain text ready to be injected into a turn.  Converters degrade
gracefully: if a preferred tool is unavailable they fall back to the next best
option, and if everything fails they return a short stub describing what to
install so the user knows what's missing.
"""
from __future__ import annotations

import csv
import html as html_lib
import io
import json
import re
import shutil
import subprocess
import tempfile
import urllib.request
from pathlib import Path
from typing import Any


# ── public dispatch ───────────────────────────────────────────────────────────

def convert(kind: str, path: Path | None, url: str = "", name: str = "") -> str:
    """Route an attachment to the right converter.  Returns extracted text."""
    label = name or (path.name if path else kind)
    try:
        return _dispatch(kind, path, url, label)
    except Exception as exc:
        return f"[{label}: converter error — {exc}]"


def _dispatch(kind: str, path: Path | None, url: str, label: str) -> str:
    if kind in ("text", "markdown"):
        return _plain_text(path, label)
    if kind == "html":
        return _html_file(path, label)
    if kind == "webpage":
        return _webpage(url or label)
    if kind == "pdf":
        return _pdf(path, label)
    if kind == "image":
        return f"[{label}: image — will be passed natively to vision model]"
    if kind == "audio file":
        return _audio_file(path, label)
    if kind == "video":
        return _video(path, label)
    if kind == "structured data":
        return _structured(path, label)
    if kind == "spreadsheet":
        return _spreadsheet(path, label)
    if kind == "document":
        return _office_doc(path, label)
    if kind == "presentation":
        return _presentation(path, label)
    if kind == "latex":
        return _plain_text(path, label)   # LaTeX source is readable as-is
    if kind == "mermaid":
        return _plain_text(path, label)   # Mermaid source is the artefact
    if kind == "ebook":
        return _epub(path, label)
    return f"[{label}: no converter for kind '{kind}']"


# ── text / markup ─────────────────────────────────────────────────────────────

def _plain_text(path: Path | None, label: str) -> str:
    if not path or not path.exists():
        return f"[{label}: file not found]"
    return path.read_text(errors="replace").strip()[:40_000]


def _html_file(path: Path | None, label: str) -> str:
    if not path or not path.exists():
        return f"[{label}: file not found]"
    return _strip_html(path.read_text(errors="replace"))


def _webpage(url: str) -> str:
    """Fetch a URL and extract readable text with trafilatura (preferred) or stdlib."""
    try:
        import trafilatura  # type: ignore
        downloaded = trafilatura.fetch_url(url)
        if downloaded:
            text = trafilatura.extract(downloaded, include_links=False,
                                       include_comments=False)
            if text:
                return text.strip()[:40_000]
    except ImportError:
        pass
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "LAWRENCE/1"})
        with urllib.request.urlopen(req, timeout=15) as r:
            raw = r.read().decode(errors="replace")
        return _strip_html(raw)
    except Exception as exc:
        return f"[webpage {url}: fetch failed — {exc}]"


def _strip_html(raw: str) -> str:
    raw = re.sub(r"<(script|style|head)[^>]*>.*?</\1>", "", raw,
                 flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<[^>]+>", " ", raw)
    text = html_lib.unescape(text)
    text = re.sub(r" {2,}", " ", re.sub(r"\n{3,}", "\n\n", text)).strip()
    return text[:40_000]


# ── PDF ───────────────────────────────────────────────────────────────────────

def _pdf(path: Path | None, label: str) -> str:
    if not path or not path.exists():
        return f"[{label}: file not found]"

    # 1. pdftotext (poppler-utils)
    if shutil.which("pdftotext"):
        try:
            r = subprocess.run(
                ["pdftotext", "-layout", str(path), "-"],
                capture_output=True, text=True, timeout=30,
            )
            if r.returncode == 0 and r.stdout.strip():
                return r.stdout.strip()[:40_000]
        except Exception:
            pass

    # 2. PyMuPDF
    try:
        import fitz  # type: ignore
        doc = fitz.open(str(path))
        pages = [p.get_text() for p in doc]
        return ("\n\n".join(pages)).strip()[:40_000]
    except ImportError:
        pass

    # 3. pdfminer.six
    try:
        from pdfminer.high_level import extract_text as pm_extract  # type: ignore
        text = pm_extract(str(path))
        if text and text.strip():
            return text.strip()[:40_000]
    except ImportError:
        pass

    return (f"[{label}: PDF text extraction unavailable — "
            "install poppler-utils (pdftotext), PyMuPDF (pip install pymupdf), "
            "or pdfminer.six (pip install pdfminer.six)]")


# ── audio / video ─────────────────────────────────────────────────────────────

def _audio_file(path: Path | None, label: str) -> str:
    """Transcribe an audio file with faster-whisper or whisper-cli."""
    if not path or not path.exists():
        return f"[{label}: file not found]"
    try:
        from faster_whisper import WhisperModel  # type: ignore
        model = WhisperModel("base", device="cpu", compute_type="int8")
        segs, _ = model.transcribe(str(path), beam_size=1)
        transcript = " ".join(s.text.strip() for s in segs).strip()
        return f"[Audio transcript — {label}]\n{transcript}" if transcript else f"[{label}: no speech detected]"
    except ImportError:
        pass
    if shutil.which("whisper-cli") or shutil.which("whisper"):
        try:
            cmd = shutil.which("whisper-cli") or "whisper"
            r = subprocess.run(
                [cmd, str(path), "--output_format", "txt", "--output_dir", str(path.parent)],
                capture_output=True, text=True, timeout=120,
            )
            out_txt = path.with_suffix(".txt")
            if out_txt.exists():
                t = out_txt.read_text(errors="replace").strip()
                out_txt.unlink(missing_ok=True)
                return f"[Audio transcript — {label}]\n{t}"
        except Exception:
            pass
    return (f"[{label}: audio transcription unavailable — "
            "install faster-whisper (pip install faster-whisper)]")


def _video(path: Path | None, label: str) -> str:
    """Extract audio from a video and transcribe it."""
    if not path or not path.exists():
        return f"[{label}: file not found]"
    if not shutil.which("ffmpeg"):
        return (f"[{label}: video audio extraction unavailable — install ffmpeg]")
    with tempfile.TemporaryDirectory(prefix="lk-vid-") as tmp:
        audio_out = Path(tmp) / "audio.wav"
        try:
            r = subprocess.run(
                ["ffmpeg", "-i", str(path), "-vn", "-acodec", "pcm_s16le",
                 "-ar", "16000", "-ac", "1", str(audio_out), "-y"],
                capture_output=True, timeout=120,
            )
            if r.returncode == 0 and audio_out.exists():
                transcript = _audio_file(audio_out, label)
                return f"[Video audio — {label}]\n{transcript}"
        except Exception as exc:
            return f"[{label}: ffmpeg error — {exc}]"
    return f"[{label}: video extraction failed]"


# ── structured data ───────────────────────────────────────────────────────────

def _structured(path: Path | None, label: str) -> str:
    if not path or not path.exists():
        return f"[{label}: file not found]"
    suffix = path.suffix.lower()
    text = path.read_text(errors="replace")

    if suffix == ".json":
        try:
            return json.dumps(json.loads(text), indent=2, default=str)[:40_000]
        except Exception:
            return text[:40_000]

    if suffix == ".jsonl":
        out = []
        for line in text.splitlines():
            line = line.strip()
            if not line:
                continue
            try:
                out.append(json.dumps(json.loads(line), default=str))
            except Exception:
                out.append(line)
            if len(out) >= 500:
                break
        return "\n".join(out)

    if suffix in (".yaml", ".yml"):
        try:
            import yaml  # type: ignore
            return json.dumps(yaml.safe_load(text), indent=2, default=str)[:40_000]
        except ImportError:
            return text[:40_000]   # return raw YAML if pyyaml missing
        except Exception:
            return text[:40_000]

    if suffix == ".xml":
        return _strip_html(text)

    return text[:40_000]


# ── spreadsheets ──────────────────────────────────────────────────────────────

def _spreadsheet(path: Path | None, label: str) -> str:
    if not path or not path.exists():
        return f"[{label}: file not found]"
    suffix = path.suffix.lower()

    if suffix in (".csv", ".tsv"):
        delim = "\t" if suffix == ".tsv" else ","
        try:
            text = path.read_text(errors="replace")
            reader = csv.reader(io.StringIO(text), delimiter=delim)
            rows = ["\t".join(row) for row in reader]
            return "\n".join(rows[:1000])
        except Exception as exc:
            return f"[{label}: csv read error — {exc}]"

    # xlsx / xlsm
    if suffix in (".xlsx", ".xlsm"):
        try:
            import openpyxl  # type: ignore
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                rows = []
                for row in sheet.iter_rows(max_row=500, values_only=True):
                    rows.append("\t".join("" if v is None else str(v) for v in row))
                parts.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
            return "\n\n".join(parts)[:40_000]
        except ImportError:
            pass

    # ODS
    if suffix == ".ods":
        try:
            import odf.opendocument  # type: ignore
            import odf.table          # type: ignore
            doc = odf.opendocument.load(str(path))
            parts = []
            for sheet in doc.spreadsheet.childNodes:
                if not hasattr(sheet, "childNodes"):
                    continue
                rows = []
                for row in sheet.childNodes:
                    cells = [str(c) for c in row.childNodes if hasattr(c, "childNodes")]
                    rows.append("\t".join(cells))
                parts.append("\n".join(rows[:500]))
            return "\n\n".join(parts)[:40_000]
        except ImportError:
            pass

    return (f"[{label}: spreadsheet converter unavailable for {suffix} — "
            "install openpyxl (pip install openpyxl) for xlsx]")


# ── office documents ──────────────────────────────────────────────────────────

def _office_doc(path: Path | None, label: str) -> str:
    if not path or not path.exists():
        return f"[{label}: file not found]"
    suffix = path.suffix.lower()

    # docx
    if suffix == ".docx":
        try:
            import docx  # type: ignore
            doc = docx.Document(str(path))
            return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())[:40_000]
        except ImportError:
            pass

    # odt / rtf / doc via pandoc
    if shutil.which("pandoc"):
        try:
            r = subprocess.run(
                ["pandoc", str(path), "-t", "plain", "--wrap=none"],
                capture_output=True, text=True, timeout=30,
            )
            if r.returncode == 0 and r.stdout.strip():
                return r.stdout.strip()[:40_000]
        except Exception:
            pass

    return (f"[{label}: office document converter unavailable for {suffix} — "
            "install python-docx (pip install python-docx) for docx, or pandoc for other formats]")


# ── presentations ─────────────────────────────────────────────────────────────

def _presentation(path: Path | None, label: str) -> str:
    if not path or not path.exists():
        return f"[{label}: file not found]"
    suffix = path.suffix.lower()

    if suffix in (".pptx", ".pptm"):
        try:
            import pptx  # type: ignore
            prs = pptx.Presentation(str(path))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    slides.append(f"[Slide {i}]\n" + "\n".join(texts))
            return "\n\n".join(slides)[:40_000]
        except ImportError:
            pass

    if shutil.which("pandoc"):
        try:
            r = subprocess.run(
                ["pandoc", str(path), "-t", "plain", "--wrap=none"],
                capture_output=True, text=True, timeout=30,
            )
            if r.returncode == 0 and r.stdout.strip():
                return r.stdout.strip()[:40_000]
        except Exception:
            pass

    return (f"[{label}: presentation converter unavailable for {suffix} — "
            "install python-pptx (pip install python-pptx)]")


# ── ebooks ────────────────────────────────────────────────────────────────────

def _epub(path: Path | None, label: str) -> str:
    if not path or not path.exists():
        return f"[{label}: file not found]"
    try:
        import ebooklib  # type: ignore
        import ebooklib.epub  # type: ignore
        book = ebooklib.epub.read_epub(str(path))
        parts = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            text = _strip_html(item.get_content().decode(errors="replace"))
            if text.strip():
                parts.append(text.strip())
        return "\n\n".join(parts)[:40_000]
    except ImportError:
        pass
    if shutil.which("pandoc"):
        try:
            r = subprocess.run(
                ["pandoc", str(path), "-t", "plain", "--wrap=none"],
                capture_output=True, text=True, timeout=30,
            )
            if r.returncode == 0 and r.stdout.strip():
                return r.stdout.strip()[:40_000]
        except Exception:
            pass
    return (f"[{label}: epub converter unavailable — "
            "install ebooklib (pip install ebooklib) or pandoc]")
